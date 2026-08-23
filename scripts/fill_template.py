"""Fill the UniHack prototype template with our content."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

BLUE = RGBColor(0x1D, 0x4E, 0xD8)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x47, 0x55, 0x69)

p = Presentation(r'D:\[EXT] UniHack-Protoype Template.pptx')
SLIDES = list(p.slides)


def add_text(slide, text, top, left=Emu(200000), width=Emu(8700000), size=12, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, Emu(3200000))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = para.add_run()
        if line.startswith("**"):
            r.text = line.strip("*")
            r.font.bold = True
            r.font.color.rgb = BLUE
            r.font.size = Pt(size + 1)
        else:
            r.text = line
            r.font.size = Pt(size)
            r.font.color.rgb = color
        para.space_after = Pt(4)


# S2 team details
for sh in SLIDES[1].shapes:
    if sh.has_text_frame and "Team name" in sh.text_frame.text:
        tf = sh.text_frame
        tf.clear()
        for j, line in enumerate(["Team Details", "", "Team name: Team Statiz",
                                  "Team leader name: Rahul G"]):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            r = para.add_run()
            r.text = line
            r.font.size = Pt(16)

# S3 brief
add_text(SLIDES[2],
    "SpecSheet \u2192 Smart Catalog turns minimal product information (part number, brand, short "
    "description, manufacturer URL, reference PDFs) into rich, governed product intelligence.\n"
    "Given sparse inputs, the pipeline scrapes manufacturer pages and reference documents, extracts "
    "spec tables from PDFs, auto-classifies each product into UNSPSC taxonomy with a confidence "
    "score, generates the six description types, up to 20 features and 50 attribute triplets "
    "(label/value/UOM), and writes the exact 252-column delivery format \u2014 with a human review "
    "console for governance.\n"
    "Result on demo data: 145 products, 100% family classification accuracy, 116 cross-supplier "
    "cross-references, 8-second CPU runtime, $0 API cost.",
    Emu(1400000), size=14)

# S4 three questions
add_text(SLIDES[3],
    "**1 \u00b7 Enrichment from minimal inputs\n"
    "Manufacturer page scraping (static + headless-browser fallback) extracts JSON-LD, meta tags and "
    "page text; reference PDFs (owner manuals, installation guides, spec sheets) are parsed into "
    "key/value tables. An LLM pass (structured JSON output) fuses all evidence into the delivery "
    "fields: 6 description variants, 20 features, 50 attribute triplets with UOMs, dimensions, "
    "certifications, images and document links. A heuristic JSON-LD path does the same offline.\n"
    "**2 \u00b7 Accuracy and trust\n"
    "- Confidence scoring on every classification (badge shown per product)\n"
    "- Multi-source verification: LLM values never overwrite hard JSON-LD facts; attributes "
    "cross-checked between page and PDF\n"
    "- Source grounding: every enriched attribute stores its source document + header\n"
    "- Rule-based validation: unit/range checks, required-attribute flags per family\n"
    "- Human review console: approve/correct before export \u2014 governed by construction\n"
    "- Anti-hallucination prompt constraints: leave-blank rules; no invented certifications\n"
    "**3 \u00b7 Enterprise scalability\n"
    "- Fuzzy column detection (RapidFuzz): new manufacturers and unseen file layouts map "
    "automatically \u2014 no code changes, no hardcoding\n"
    "- Embedding classification scales to full UNSPSC/eCl@ss (41k classes) by swapping taxonomy "
    "data \u2014 architecture unchanged\n"
    "- Row-by-row streaming + disk-cached fetches handle large catalogs and cheap re-runs\n"
    "- Continuous updates: re-run any row; review state persists",
    Emu(2100000), size=10)

# S5 opportunities
add_text(SLIDES[4],
    "**How different from existing ideas?\n"
    "No mainstream PIM/search vendor natively auto-classifies to industrial taxonomies "
    "(UNSPSC/eCl@ss/ETIM); only one (inRiver) partially extracts datasheets; enrichment leaders "
    "enrich their own search index, not an exportable governed catalog. Partly proved "
    "cross-referencing data is worth $500M \u2014 but automotive-only. We combine all of it, "
    "self-serve.\n"
    "**How it solves the problem statement\n"
    "Minimal product information (part no, brand, short desc) + manufacturer URLs is exactly our "
    "input; the 252-column delivery format is exactly our output \u2014 populated dynamically with "
    "unseen evaluation data handled by fuzzy schema detection.\n"
    "**USP\n"
    "1. $0 cost, CPU-only, offline-capable (graceful LLM fallback) \u2014 vs $45\u2013100K/yr tools\n"
    "2. Every AI value carries evidence: confidence + source span\n"
    "3. Cross-supplier substitute matching built in\n"
    "4. 8-second end-to-end runtime on a laptop",
    Emu(2350000), size=10)

# S6 features
add_text(SLIDES[5],
    "- Fuzzy schema ingestion: any supplier's CSV/XLSX column names auto-mapped (in\u2192mm, comma decimals)\n"
    "- Manufacturer URL scraping with headless-browser fallback \u00b7 PDF spec-sheet extraction\n"
    "- UNSPSC auto-classification with confidence scores and top-5 candidate evidence\n"
    "- LLM enrichment: 6 description types, 20 features, 50 attribute triplets (label/value/UOM)\n"
    "- Cross-supplier part cross-referencing + substitute search API\n"
    "- Human review console: approve/correct, confidence badges, source-tagged attributes\n"
    "- Delivery export: all 252 static headers populated exactly (CSV + XLSX)\n"
    "- Multi-provider free LLM layer (Cerebras/Groq/Gemini) with deterministic offline fallback",
    Emu(1500000), size=13)

# S7 flow diagram
SLIDES[6].shapes.add_picture("demo/assets/flow.png", Emu(150000), Emu(1450000), height=Emu(3300000))

# S8 wireframe = console
SLIDES[7].shapes.add_picture("demo/assets/mvp_console.png", Emu(150000), Emu(1400000), height=Emu(3300000))

# S9 architecture
SLIDES[8].shapes.add_picture("demo/assets/arch.png", Emu(150000), Emu(1500000), height=Emu(3200000))

# S10 tech
add_text(SLIDES[9],
    "Python 3.14 \u00b7 FastAPI + Uvicorn (API & review console) \u00b7 sentence-transformers MiniLM "
    "(local embeddings)\n"
    "PyMuPDF + pdfplumber (PDF extraction) \u00b7 RapidFuzz (fuzzy matching) \u00b7 pandas/openpyxl (I/O)\n"
    "Playwright/Chromium (JS-page scraping) \u00b7 httpx (cached fetch layer)\n"
    "LLM layer: Cerebras \u2192 Groq \u2192 Gemini free tiers, OpenAI-compatible JSON mode",
    Emu(2100000), size=14)

# S11 cost
add_text(SLIDES[10],
    "Prototype: $0.\n"
    "- LLMs: free tiers (Cerebras ~30K TPM / Groq / Gemini), offline heuristic fallback\n"
    "- Embeddings: local on CPU \u2014 no API cost\n"
    "- Hosting: Cloudflare quick tunnel (demo); Render free tier for a persistent MVP\n"
    "At scale: paid LLM tier \u2248 $0.001\u20130.005 per product; 100K products \u2248 $100\u2013500 one-time.",
    Emu(1500000), size=14)

# S12 snapshots (two images side by side)
SLIDES[11].shapes.add_picture("demo/assets/mvp_console.png", Emu(120000), Emu(1400000), height=Emu(2100000))
SLIDES[11].shapes.add_picture("demo/assets/mvp_api.png", Emu(4750000), Emu(1400000), height=Emu(2100000))
add_text(SLIDES[11], "Review console (left) \u00b7 substitute search API (right)",
         Emu(3700000), left=Emu(3600000), size=11, color=GRAY)

# S13 future
add_text(SLIDES[12],
    "- Engineering-constraint substitution validation (fit/load margins, not just dimensions)\n"
    "- OCR path for scanned datasheets\n"
    "- Full eCl@ss/ETIM taxonomy trees + multi-language descriptions\n"
    "- Image enrichment: product-photo attribute extraction\n"
    "- Incremental catalog sync + enterprise SSO for review workflow",
    Emu(1550000), size=14)

# S14 links
add_text(SLIDES[13],
    "GitHub Public Repository:\nhttps://github.com/Rahulxoxo69/Statiz\n\n"
    "Demo Video:\nhttps://github.com/Rahulxoxo69/Statiz/blob/main/demo/demo_video.mp4\n\n"
    "Working Prototype:\nhttps://day-atm-become-bureau.trycloudflare.com\n"
    "(live review console \u2014 search, filter, approve, export CSV; substitutes API at /api/substitutes/BRG-BEA0000)",
    Emu(1600000), size=14)

# remove guidelines slide 1
xml_slides = p.slides._sldIdLst
xml_slides.remove(list(xml_slides)[0])

p.save(r'D:\projects1\hackathon-industrial\demo\UniHack_Prototype.pptx')
print("saved demo/UniHack_Prototype.pptx")
