"""Update text content in the user's final-formatted deck without touching layout.
Only rewrites text inside existing TextBox shapes; fonts/sizes/positions preserved.
"""
from pptx import Presentation
from pptx.util import Pt
import copy

SRC = r"D:\UniHack_Prototype_final.pptx"
OUT = r"D:\projects1\hackathon-industrial\demo\UniHack_Prototype_final.pptx"

p = Presentation(SRC)
slides = list(p.slides)


def set_text(shape, lines, size=None):
    """Replace a shape's text, reusing the first run's formatting for all lines."""
    tf = shape.text_frame
    # capture template font from first run
    tmpl = None
    for para in tf.paragraphs:
        for r in para.runs:
            tmpl = copy.deepcopy(r.font)
            break
        if tmpl is not None:
            break
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = para.add_run()
        r.text = line
        if tmpl is not None:
            r.font._rPr.set("sz", tmpl._rPr.get("sz")) if tmpl._rPr.get("sz") else None
            try:
                r.font.bold = tmpl.bold
                r.font.color.rgb = tmpl.color.rgb
                r.font.name = tmpl.name
            except Exception:
                pass
        if size:
            r.font.size = Pt(size)
    return shape


def find(slide, startswith, contains=None):
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t.startswith(startswith) and (contains is None or contains in t):
                return sh
    return None


# ---- slide 2: Brief — lead with real dataset, new numbers
sh = find(slides[1], "Statiz turns minimal")
set_text(sh, [
    "Statiz turns minimal product information — part number, brand, short description (and, when available, manufacturer URLs and reference PDFs) — into rich, governed product intelligence delivered in the exact 252-column expected output format.",
    "On the provided sample dataset (1,000 rows, six sparse input columns), the rule + AI pipeline resolves trading brands from manufacturer names, classifies every row into Dept/Class/Fine with UNSPSC codes, and extracts technical specs (dimensions, grit, voltage, pack counts, materials) into typed attribute triplets — generating all six description variants, features and attributes dynamically. Nothing is hardcoded to the sample: fuzzy column detection and generic parsing rules handle unseen evaluation data.",
    "Results: 1,000/1,000 rows classified (zero empty), 640 with UNSPSC codes, 627 with extracted attributes; URL-scraping mode (JSON-LD + PDF datasheets) adds deeper enrichment when links exist; 8-second CPU runtime at $0 API cost, with a human review console for governance.",
], size=12)

# ---- slide 3: enrichment — description-driven first
sh = find(slides[2], "Manufacturer page scraping")
set_text(sh, [
    "Description-driven enrichment (no URLs required): a rule-based parser resolves the trading brand (Freud Inc \u2192 Diablo; tokens like 3M/Milwaukee found in the description win over distributor names), classifies product type via a 60+ pattern taxonomy (abrasives, lighting, tools, lumber, fasteners, appliances, safety\u2026) into Dept/Class/Fine + UNSPSC, and extracts specs with typed regexes \u2014 dimensions, fractions, grit (P150), voltage, amps, Kelvin, lumens, dBA, pack counts, materials (SST \u2192 Stainless Steel), colors, teeth \u2014 each becoming a label/value/UOM attribute triplet.",
    "URL mode (when links exist): manufacturer pages scraped via static fetch + headless-browser fallback (JSON-LD, meta, page text); reference PDFs (manuals, spec sheets) parsed with PyMuPDF + pdfplumber into key/value tables; a free-tier LLM pass (structured JSON) fuses all evidence into the full delivery schema.",
    "A deterministic heuristic path does the same offline \u2014 the system never fails for lack of an API key.",
], size=11)

# ---- slide 4: accuracy — add concrete rule examples
sh = find(slides[3], "Confidence scoring on every classification")
set_text(sh, [
    "Confidence scoring on every classification, shown as a badge per product (embedding similarity + LLM self-assessed confidence).",
    "Multi-source verification: LLM values never overwrite hard JSON-LD facts; attributes cross-checked between page text and PDF; brand resolved from multiple signals with a deterministic priority order.",
    "Rule-based validation: typed spec regexes reject malformed values; unit/range checks; required-attribute flags per product family; every enriched attribute stores its source (PDF name / header / regex on description).",
    "Human review console: approve/correct each row before export \u2014 governed by construction.",
    "Anti-hallucination constraints: leave-blank rules for unknown fields; no invented certifications or numbers.",
], size=11)

# ---- slide 5: scalability — cite the 1000-row run
sh = find(slides[4], "Fuzzy column detection (RapidFuzz)")
set_text(sh, [
    "Fuzzy column detection (RapidFuzz): new manufacturers and unseen file layouts map automatically \u2014 no code changes, no hardcoding. Proven live: the solver processed the 1,000-row sample dataset with different field combinations and zero per-row configuration.",
    "Embedding classification scales to the full UNSPSC/eCl@ss tree (41k classes) by swapping taxonomy data \u2014 architecture unchanged. Rule tables are data-driven lists, extensible in minutes.",
    "Row-by-row streaming with disk-cached fetches handles large catalogs and cheap re-runs; 1,000 rows enrich in seconds on a laptop CPU.",
    "Continuous updates: re-run any row; review state persists across runs.",
], size=11)

# ---- slide 7: features — add description-driven bullet
sh = find(slides[6], "Fuzzy schema ingestion")
set_text(sh, [
    "Description-driven enrichment from minimal inputs: brand resolution, 60+ product-type taxonomy rules, typed spec extraction (label/value/UOM triplets)",
    "Fuzzy schema ingestion: any supplier's CSV/XLSX column names auto-mapped (in\u2192mm, comma decimals)",
    "Manufacturer URL scraping with headless-browser fallback \u00b7 PDF spec-sheet extraction (PyMuPDF + pdfplumber)",
    "UNSPSC auto-classification with confidence scores and top-5 candidate evidence",
    "LLM enrichment (free tiers): 6 description types, 20 features, 50 attribute triplets, dims, certifications",
    "Cross-supplier part cross-referencing + substitute search API",
    "Human review console: approve/correct, confidence badges, source-tagged attributes",
    "Delivery export: all 252 static headers populated exactly, CSV + XLSX, fully dynamic on unseen data",
], size=12)

# ---- slide 12: cost unchanged; slide 15 links unchanged; slide 14 future unchanged

p.save(OUT)
print("saved", OUT)
